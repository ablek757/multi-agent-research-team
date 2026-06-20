"""PPT 格式器：使用 python-pptx 生成幻灯片。"""

import re
from pathlib import Path
from typing import List

from agent.config import Config
from agent.output.base import OutputArtifact, OutputFormatter
from agent.research_state import ResearchState


class SlidesFormatter(OutputFormatter):
    """将研究报告转换为 PPT 幻灯片。"""

    name = "slides"
    mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    extension = ".pptx"

    def __init__(self, config: Config, style_profile=None):
        super().__init__(config, style_profile=style_profile)
        self.template = config.output.slides_template

    def format(self, topic: str, state: ResearchState) -> OutputArtifact:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError as exc:
            raise ImportError(
                "python-pptx is required for slides output. Install with: pip install python-pptx"
            ) from exc

        report_body = self._build_report_body(topic, state)
        prs = Presentation()

        # 标题页
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = topic
        subtitle = title_slide.placeholders[1]
        subtitle.text = self.config.report.title

        # 将报告正文按二级标题拆分
        sections = self._split_sections(report_body)
        for heading, bullets in sections:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = heading[:60]
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for bullet in bullets[:8]:
                p = tf.add_paragraph()
                p.text = bullet[:120]
                p.level = 0
                p.font.size = Pt(18)

        # 来源页
        if state.sources:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "参考来源"
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for s in state.sources[:10]:
                p = tf.add_paragraph()
                p.text = f"[{s.index}] {s.title[:80]}"
                p.level = 0
                p.font.size = Pt(14)

        # 保存到内存字符串/bytes？python-pptx 必须写文件
        file_path = str(Path(self.config.report.output_dir) / f"{self._safe_name(topic)}.pptx")
        Path(self.config.report.output_dir).mkdir(parents=True, exist_ok=True)
        prs.save(file_path)

        return OutputArtifact(
            format=self.name,
            file_path=file_path,
            mime_type=self.mime_type,
            metadata={"slide_count": len(prs.slides), "source_count": len(state.sources)},
        )

    def _split_sections(self, text: str) -> List[tuple]:
        """按 Markdown 二级标题拆分章节。"""
        sections = []
        current_heading = "要点"
        current_bullets: List[str] = []

        for line in text.splitlines():
            line = line.strip()
            if not line:
                continue
            if line.startswith("## "):
                if current_bullets:
                    sections.append((current_heading, current_bullets))
                current_heading = line.lstrip("# ").strip()
                current_bullets = []
            elif line.startswith("- ") or line.startswith("* "):
                current_bullets.append(line.lstrip("- *").strip())
            elif line and not line.startswith("#"):
                # 将段落按句拆为 bullet
                sentences = re.split(r"(?<=[。！？.!?])\s+", line)
                current_bullets.extend(sentences)

        if current_bullets:
            sections.append((current_heading, current_bullets))

        if not sections:
            sections.append(("研究摘要", [text[:500]]))
        return sections

    def _safe_name(self, topic: str) -> str:
        return re.sub(r'[\\/:*?"<>|]+', "_", topic)[:50]
